#!/usr/bin/env python3
"""Generate PPTX from HTML presentation content."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import math

# ─── Colors ───
GREEN_DEEP = RGBColor(0x1A, 0x3A, 0x2A)
GREEN_DARK = RGBColor(0x2C, 0x5F, 0x2D)
GREEN_MID = RGBColor(0x3D, 0x7A, 0x4A)
GREEN_LIGHT = RGBColor(0x5F, 0xA8, 0x6B)
GREEN_PALE = RGBColor(0xA8, 0xD5, 0xB0)
CREAM = RGBColor(0xF5, 0xF0, 0xE8)
CREAM_LIGHT = RGBColor(0xFA, 0xF7, 0xF2)
GOLD = RGBColor(0xC4, 0xA3, 0x5A)
GOLD_LIGHT = RGBColor(0xE8, 0xD5, 0xA0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_DARK = RGBColor(0x1A, 0x2A, 0x1E)
TEXT_MID = RGBColor(0x3A, 0x4A, 0x3E)
BARK = RGBColor(0x5C, 0x4A, 0x3A)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_gradient_bg(slide, c1, c2):
    """Approximate gradient with solid dark green."""
    add_bg(slide, c1)

def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, italic=False, alignment=PP_ALIGN.LEFT,
                 font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_accent_bar(slide, color=GREEN_LIGHT):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.06), H)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def add_section_tag(slide, text, top=Inches(0.5), color=None, left=Inches(0.9)):
    c = color if color else RGBColor(0x88, 0xAA, 0x88)
    add_text_box(slide, left, top, Inches(8), Inches(0.35), text,
                 font_size=11, color=c, bold=True)

def add_slide_title(slide, text, top=Inches(0.8), color=None, left=Inches(0.9)):
    c = color if color else GREEN_DEEP
    add_text_box(slide, left, top, Inches(11), Inches(0.8), text,
                 font_size=36, color=c, bold=True, font_name='Georgia')

def add_table(slide, left, top, width, rows_data, col_widths=None, header_color=GREEN_DARK):
    rows = len(rows_data)
    cols = len(rows_data[0])
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, Inches(0.45 * rows))
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for r, row_data in enumerate(rows_data):
        for c, cell_text in enumerate(row_data):
            cell = table.cell(r, c)
            cell.text = str(cell_text)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(13)
                paragraph.font.name = 'Calibri'
                if r == 0:
                    paragraph.font.color.rgb = WHITE
                    paragraph.font.bold = True
                    paragraph.font.size = Pt(11)
                else:
                    paragraph.font.color.rgb = TEXT_DARK
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color

    return table_shape

def add_bullet_list(slide, left, top, width, items, color=TEXT_DARK, font_size=15):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(len(items) * 0.45))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Calibri'
        p.space_after = Pt(6)
        p.level = 0
    return txBox

def add_numbered_list(slide, left, top, width, items, color=TEXT_DARK, num_color=GREEN_DARK, font_size=15):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(len(items) * 0.45))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = f"{i+1}.  "
        run.font.size = Pt(font_size)
        run.font.color.rgb = num_color
        run.font.bold = True
        run.font.name = 'Calibri'
        run2 = p.add_run()
        run2.text = item
        run2.font.size = Pt(font_size)
        run2.font.color.rgb = color
        run2.font.name = 'Calibri'
        p.space_after = Pt(6)
    return txBox

def add_stat_card(slide, left, top, width, height, number, label,
                  bg_color=None, num_color=GREEN_DARK, label_color=TEXT_MID):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    if bg_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = WHITE
    shape.line.fill.background()
    shape.shadow.inherit = False

    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = str(number)
    run.font.size = Pt(28)
    run.font.color.rgb = num_color
    run.font.bold = True
    run.font.name = 'Georgia'

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = label
    run2.font.size = Pt(10)
    run2.font.color.rgb = label_color
    run2.font.bold = True
    run2.font.name = 'Calibri'

def add_card(slide, left, top, width, height, title, text, bg_color=WHITE,
             title_color=GREEN_DARK, text_color=TEXT_MID, icon=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.15)

    p = tf.paragraphs[0]
    if icon:
        run_icon = p.add_run()
        run_icon.text = icon + "  "
        run_icon.font.size = Pt(20)
    run = p.add_run()
    run.text = title
    run.font.size = Pt(16)
    run.font.color.rgb = title_color
    run.font.bold = True
    run.font.name = 'Georgia'

    p2 = tf.add_paragraph()
    run2 = p2.add_run()
    run2.text = text
    run2.font.size = Pt(12)
    run2.font.color.rgb = text_color
    run2.font.name = 'Calibri'
    p2.space_before = Pt(6)


def separator_slide(title, subtitle, icon=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, GREEN_DEEP)
    if icon:
        add_text_box(slide, Inches(0), Inches(2.2), W, Inches(0.8), icon,
                     font_size=48, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0), Inches(3.0), W, Inches(1.0), title,
                 font_size=44, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER, font_name='Georgia')
    # Gold line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(5.8), Inches(4.1), Inches(1.7), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = GOLD
    line.line.fill.background()
    add_text_box(slide, Inches(0), Inches(4.3), W, Inches(0.5), subtitle,
                 font_size=16, color=RGBColor(0x99, 0xBB, 0x99), alignment=PP_ALIGN.CENTER)
    return slide

def dark_content_slide(title, section_tag=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, GREEN_DEEP)
    add_accent_bar(slide, GREEN_LIGHT)
    if section_tag:
        add_section_tag(slide, section_tag.upper(), color=RGBColor(0x77, 0x99, 0x77))
    add_slide_title(slide, title, color=WHITE)
    return slide

def cream_content_slide(title, section_tag=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, CREAM_LIGHT)
    add_accent_bar(slide, GREEN_DARK)
    if section_tag:
        add_section_tag(slide, section_tag.upper(), color=RGBColor(0x88, 0x99, 0x88), left=Inches(0.9))
    add_slide_title(slide, title, color=GREEN_DEEP)
    return slide


# ════════════════════════════════════════════════════════════════
# SLIDE 1: PORTADA
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, GREEN_DEEP)

add_text_box(slide, Inches(0), Inches(1.5), W, Inches(0.6), "🍃",
             font_size=48, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(2.2), Inches(10.3), Inches(1.4),
             "Obtención de Repelente Natural\na partir de Aceite Esencial de Eucalipto",
             font_size=38, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER, font_name='Georgia')
add_text_box(slide, Inches(0), Inches(3.7), W, Inches(0.5),
             "Proyecto Final — Ingeniería Química",
             font_size=20, color=GOLD_LIGHT, italic=True, alignment=PP_ALIGN.CENTER, font_name='Georgia')

line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.8), Inches(4.4), Inches(1.7), Inches(0.02))
line.fill.solid(); line.fill.fore_color.rgb = GOLD; line.line.fill.background()

add_text_box(slide, Inches(0), Inches(4.6), W, Inches(0.4),
             "BRESSO, Romina  ·  MACAGNO, Micaela",
             font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(0), Inches(5.1), W, Inches(0.3),
             "Directora: Ing. Sposetti, Patricia  |  Tutora: Ing. Dra. Taverna, María Eugenia",
             font_size=12, color=RGBColor(0x99,0x99,0x99), alignment=PP_ALIGN.CENTER)

line2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.8), Inches(5.6), Inches(1.7), Inches(0.015))
line2.fill.solid(); line2.fill.fore_color.rgb = GOLD; line2.line.fill.background()

add_text_box(slide, Inches(0), Inches(5.8), W, Inches(0.3),
             "Universidad Tecnológica Nacional — Facultad Regional San Francisco",
             font_size=12, color=RGBColor(0x99,0x99,0x99), alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(0), Inches(6.2), W, Inches(0.3), "2025",
             font_size=12, color=RGBColor(0x77,0x77,0x77), alignment=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDE 2: OBJETIVOS
# ════════════════════════════════════════════════════════════════
slide = cream_content_slide("Objetivos del Proyecto", "Bloque 1 · Portada y Contexto")
objectives = [
    "Describir la materia prima y el producto final",
    "Identificar envases y rotulación",
    "Definir localización de la planta",
    "Detallar etapas del proceso productivo",
    "Garantizar calidad mediante análisis",
    "Establecer condiciones de seguridad e higiene",
    "Evaluar impacto ambiental y social",
    "Indicar leyes y normativas aplicables",
    "Estimar costo unitario, precio de venta y factibilidad",
]
add_numbered_list(slide, Inches(1.0), Inches(1.8), Inches(10), objectives)

# ════════════════════════════════════════════════════════════════
# SLIDE 3: CONTEXTO
# ════════════════════════════════════════════════════════════════
slide = cream_content_slide("Repelente Natural — Contexto", "Contexto")
cards_data = [
    ("🌍", "Cambio Global hacia lo Natural", "Creciente conciencia sobre el deterioro ambiental impulsa la búsqueda de alternativas sustentables."),
    ("⚠️", "Riesgos de Repelentes Sintéticos", "Irritación cutánea y toxicidad asociada al uso prolongado de compuestos como DEET."),
    ("♻️", "Alternativa Sustentable", "Producto biodegradable elaborado con recursos renovables y bajo impacto ambiental."),
    ("🌿", "Potencial del Eucalipto", "Propiedades repelentes comprobadas gracias al aceite esencial rico en eucaliptol."),
]
for i, (icon, title, text) in enumerate(cards_data):
    col = i % 2
    row = i // 2
    x = Inches(1.0) + col * Inches(5.8)
    y = Inches(1.9) + row * Inches(2.2)
    add_card(slide, x, y, Inches(5.4), Inches(1.8), title, text, icon=icon)

# ════════════════════════════════════════════════════════════════
# SLIDE 4: POR QUÉ
# ════════════════════════════════════════════════════════════════
slide = dark_content_slide("¿Por qué un Repelente Natural?", "Contexto")
cards_dark = [
    ("🦟", "Amenaza Sanitaria", "Los mosquitos transmiten dengue, malaria y zika. La protección es una necesidad de salud pública."),
    ("🧬", "Compuestos Activos", "El aceite de eucalipto contiene eucamalol y 4-isopropilbencil alcohol con acción repelente demostrada."),
    ("🌱", "Biodegradable", "Los productos naturales se degradan sin acumularse en el ambiente, a diferencia de los sintéticos."),
]
for i, (icon, title, text) in enumerate(cards_dark):
    x = Inches(1.0) + i * Inches(3.9)
    add_card(slide, x, Inches(1.9), Inches(3.5), Inches(2.2), title, text,
             bg_color=RGBColor(0x22, 0x44, 0x33), title_color=GREEN_PALE,
             text_color=RGBColor(0xBB, 0xCC, 0xBB), icon=icon)

# ════════════════════════════════════════════════════════════════
# SLIDE 5: TENDENCIAS
# ════════════════════════════════════════════════════════════════
slide = cream_content_slide("Tendencias Mundiales", "Contexto")
tendencias = [
    "Preocupación creciente por la salud y el bienestar",
    "Crecimiento sostenido del mercado de cosméticos naturales",
    "Mayor empleo de derivados de plantas medicinales y aromáticas",
    "Regulaciones gubernamentales a favor de productos naturales",
    "Demanda creciente de productos con bajo impacto toxicológico",
]
add_bullet_list(slide, Inches(1.0), Inches(1.9), Inches(10), tendencias, font_size=17)

# ════════════════════════════════════════════════════════════════
# SLIDE 6: SEPARADOR MERCADO
# ════════════════════════════════════════════════════════════════
separator_slide("Estudio de Mercado", "Análisis de oferta, demanda y comercialización", "📊")

# ════════════════════════════════════════════════════════════════
# SLIDE 7: TENDENCIAS MERCADO
# ════════════════════════════════════════════════════════════════
slide = dark_content_slide("Tendencias del Mercado", "Estudio de Mercado")
stats = [("USD 4.150 M", "Mercado mundial 2020"), ("6,85%", "CAGR 2021–2026"), ("+Info", "Consumidores informados")]
for i, (num, label) in enumerate(stats):
    x = Inches(1.0) + i * Inches(3.9)
    add_stat_card(slide, x, Inches(2.0), Inches(3.5), Inches(1.8), num, label,
                  bg_color=RGBColor(0x22,0x44,0x33), num_color=GREEN_PALE,
                  label_color=RGBColor(0x88,0xAA,0x88))

# ════════════════════════════════════════════════════════════════
# SLIDE 8: NATURAL vs SINTETICO
# ════════════════════════════════════════════════════════════════
slide = cream_content_slide("Natural vs Sintético", "Estudio de Mercado")
table_data = [
    ["Característica", "🌿 Natural", "🧪 Sintético"],
    ["Origen", "Vegetal", "Químico"],
    ["Duración", "Menor", "Mayor"],
    ["Efectividad", "Media–Alta", "Alta"],
    ["Riesgo para la salud", "Bajo", "Puede causar irritación"],
    ["Impacto ambiental", "Biodegradable", "Persistente"],
    ["Costo", "Competitivo", "Variable"],
]
add_table(slide, Inches(1.0), Inches(1.9), Inches(11), table_data)

# ════════════════════════════════════════════════════════════════
# SLIDE 9: DEMANDA
# ════════════════════════════════════════════════════════════════
slide = dark_content_slide("Proyección de la Demanda", "Estudio de Mercado")
add_text_box(slide, Inches(1.0), Inches(1.7), Inches(5), Inches(0.4),
             "Datos de Demanda (miles de L)", font_size=18, color=GREEN_PALE, bold=True, font_name='Georgia')
dem_data = [
    ["Año", "Optimista", "Pesimista"],
    ["2023", "91,85", "72,16"],
    ["2024", "100,39", "80,69"],
    ["2025", "108,93", "89,23"],
    ["2026", "117,46", "97,76"],
    ["2027", "125,99", "106,30"],
]
add_table(slide, Inches(1.0), Inches(2.2), Inches(5.5), dem_data, header_color=RGBColor(0x22,0x44,0x33))
add_text_box(slide, Inches(1.0), Inches(5.5), Inches(4), Inches(0.3),
             "Fuente: INDEC", font_size=10, color=RGBColor(0x77,0x99,0x77))

# ════════════════════════════════════════════════════════════════
# SLIDE 10: OFERTA
# ════════════════════════════════════════════════════════════════
slide = cream_content_slide("Proyección de la Oferta", "Estudio de Mercado")
of_data = [
    ["Año", "Optimista", "Pesimista"],
    ["2023", "534,70", "408,84"],
    ["2024", "582,81", "456,95"],
    ["2025", "630,92", "505,07"],
    ["2026", "679,03", "553,18"],
    ["2027", "727,15", "601,29"],
]
add_table(slide, Inches(1.0), Inches(1.9), Inches(11), of_data)

# ════════════════════════════════════════════════════════════════
# SLIDE 11: DPI Y PRODUCCION
# ════════════════════════════════════════════════════════════════
slide = dark_content_slide("Cuantificación DPI y Producción", "Estudio de Mercado")
stats11 = [("148.500", "Botellas / año (100 mL)"), ("495", "Botellas / día"), ("2 × 8h", "Turnos de trabajo")]
for i, (num, label) in enumerate(stats11):
    x = Inches(1.0) + i * Inches(3.9)
    add_stat_card(slide, x, Inches(2.5), Inches(3.5), Inches(1.8), num, label,
                  bg_color=RGBColor(0x22,0x44,0x33), num_color=GREEN_PALE,
                  label_color=RGBColor(0x88,0xAA,0x88))
add_text_box(slide, Inches(1.0), Inches(1.7), Inches(10), Inches(0.5),
             "DPI promedio: 99 mil L/año  →  Absorción: 15%  →  Producción: 49,50 L/día",
             font_size=16, color=GREEN_PALE)

# ════════════════════════════════════════════════════════════════
# SLIDE 12: PRECIOS
# ════════════════════════════════════════════════════════════════
slide = cream_content_slide("Precios y Comercialización", "Estudio de Mercado")
price_data = [
    ["Marca", "Precio"],
    ["Exilet", "$3.861"],
    ["Just", "$3.500"],
    ["OFF", "$2.143"],
    ["Farmacity", "$1.230"],
]
add_table(slide, Inches(1.0), Inches(1.9), Inches(5), price_data)

add_text_box(slide, Inches(7.0), Inches(1.7), Inches(5), Inches(0.4),
             "Canales de Distribución", font_size=18, color=GREEN_DARK, bold=True, font_name='Georgia')
canales = ["Venta al por mayor", "Distribución intensiva", "Agentes comerciales", "E-commerce"]
add_bullet_list(slide, Inches(7.0), Inches(2.2), Inches(5), canales)

add_text_box(slide, Inches(7.0), Inches(4.0), Inches(5), Inches(0.4),
             "Estrategia de Lanzamiento", font_size=18, color=GREEN_DARK, bold=True, font_name='Georgia')
estrategia = ["Soft launch en Santa Fe", "Redes sociales y marketing digital", "Precios promocionales iniciales"]
add_bullet_list(slide, Inches(7.0), Inches(4.5), Inches(5), estrategia)

# Precio promedio box
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(4.8), Inches(5), Inches(0.8))
shape.fill.solid(); shape.fill.fore_color.rgb = GREEN_DARK; shape.line.fill.background()
tf = shape.text_frame; tf.margin_left = Inches(0.2)
p = tf.paragraphs[0]; r = p.add_run()
r.text = "Precio promedio fábrica: $1.878,45/envase"; r.font.size = Pt(15)
r.font.color.rgb = WHITE; r.font.bold = True; r.font.name = 'Calibri'

# ════════════════════════════════════════════════════════════════
# SLIDE 13: SEP MATERIAS PRIMAS
# ════════════════════════════════════════════════════════════════
separator_slide("Materias Primas", "Eucalipto, solventes y aditivos", "🌿")

# ════════════════════════════════════════════════════════════════
# SLIDE 14: EUCALIPTO
# ════════════════════════════════════════════════════════════════
slide = cream_content_slide("Eucalipto Camaldulensis", "Materias Primas")
items14 = [
    "Árbol de gran porte, hasta 30 m de altura",
    "Corteza lisa, grisácea; copa extendida con ramas péndulas",
    "Hojas de 25 cm × 4 cm con cavidades de aceites esenciales",
    "Extendido en Buenos Aires, Santa Fe, Jujuy",
    "Especie seleccionada por su disponibilidad regional",
]
add_bullet_list(slide, Inches(1.0), Inches(1.9), Inches(5.5), items14)

tax_data = [
    ["Clasificación", ""],
    ["Reino", "Plantae"],
    ["División", "Magnoliophyta"],
    ["Familia", "Myrtaceae"],
    ["Género", "Eucalyptus"],
    ["Especie", "E. camaldulensis"],
]
add_table(slide, Inches(7.5), Inches(1.9), Inches(4.5), tax_data)

# ════════════════════════════════════════════════════════════════
# SLIDE 15: ACEITE ESENCIAL
# ════════════════════════════════════════════════════════════════
slide = dark_content_slide("Composición del Aceite Esencial", "Materias Primas")
ae_data = [
    ["Propiedad", "Valor"],
    ["Fórmula", "C₁₀H₁₈O"],
    ["Masa molar", "154,25 g/mol"],
    ["Densidad", "0,92 g/cm³"],
    ["P. ebullición", "176,50 °C"],
    ["Rendimiento", "2,15%"],
]
add_table(slide, Inches(1.0), Inches(2.2), Inches(5), ae_data, header_color=RGBColor(0x22,0x44,0x33))
add_text_box(slide, Inches(7.0), Inches(1.7), Inches(5), Inches(0.4),
             "Otros Componentes", font_size=18, color=GREEN_PALE, bold=True, font_name='Georgia')
add_text_box(slide, Inches(7.0), Inches(2.3), Inches(5), Inches(0.5),
             "α-pineno  ·  Limoneno  ·  Terpineol  ·  Flavonoides  ·  Taninos",
             font_size=14, color=RGBColor(0xAA,0xCC,0xAA))

shape15 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.0), Inches(3.5), Inches(5), Inches(1.2))
shape15.fill.solid(); shape15.fill.fore_color.rgb = RGBColor(0x1F,0x35,0x28); shape15.line.fill.background()
tf15 = shape15.text_frame; tf15.margin_left = Inches(0.2); tf15.margin_top = Inches(0.1); tf15.word_wrap = True
p15 = tf15.paragraphs[0]; r15a = p15.add_run()
r15a.text = "Método de extracción: "; r15a.font.size = Pt(14); r15a.font.color.rgb = WHITE; r15a.font.bold = True
r15b = p15.add_run()
r15b.text = "Lixiviación con etanol\nRendimiento superior (3%) frente a destilación por arrastre."
r15b.font.size = Pt(13); r15b.font.color.rgb = RGBColor(0xAA,0xBB,0xAA)

# ════════════════════════════════════════════════════════════════
# SLIDE 16: SOLVENTES
# ════════════════════════════════════════════════════════════════
slide = cream_content_slide("Solventes: Etanol e Isopropanol", "Materias Primas")
add_card(slide, Inches(1.0), Inches(1.9), Inches(5.4), Inches(2.5),
         "Etanol 96%",
         "Solvente principal para extracción. Alta volatilidad, secado rápido.\nDensidad: 797 kg/m³. Obtenido por fermentación.",
         icon="🧪")
add_card(slide, Inches(6.8), Inches(1.9), Inches(5.4), Inches(2.5),
         "Isopropanol",
         "Co-solvente. Densidad: 0,80 g/cm³. P. ebullición: 82 °C.\nVehículo del repelente, diluye el aceite esencial.",
         icon="⚗️")

# ════════════════════════════════════════════════════════════════
# SLIDE 17: AGUA Y GLICERINA
# ════════════════════════════════════════════════════════════════
slide = dark_content_slide("Agua Destilada y Glicerina", "Materias Primas")
add_card(slide, Inches(1.0), Inches(1.9), Inches(5.4), Inches(2.8),
         "Agua Destilada",
         "Libre de contaminantes. Líquido transparente, incoloro, inodoro.\nContenido orgánico total < 0,50 mg/L.\nSe emplea para preparar solución de isopropanol.",
         bg_color=RGBColor(0x22,0x44,0x33), title_color=GREEN_PALE, text_color=RGBColor(0xBB,0xCC,0xBB), icon="💧")
add_card(slide, Inches(6.8), Inches(1.9), Inches(5.4), Inches(2.8),
         "Glicerina (1,2,3-Propanotriol)",
         "Contenido ≥ 99% de C₃H₈O₃.\nHumectante: reduce deshidratación de la piel.\nMejora textura y sensación de aplicación.\nMiscible con agua y alcohol.",
         bg_color=RGBColor(0x22,0x44,0x33), title_color=GREEN_PALE, text_color=RGBColor(0xBB,0xCC,0xBB), icon="🫧")

# ════════════════════════════════════════════════════════════════
# SLIDE 18: PRODUCTO
# ════════════════════════════════════════════════════════════════
slide = cream_content_slide("Características del Producto", "Producto Elaborado")
fq_data = [["Parámetro","Valor"],["Densidad","0,92 g/cm³"],["Viscosidad","0,002 Pa·s"],["Saturación","20%"],["pH","4,50"]]
add_text_box(slide, Inches(1.0), Inches(1.7), Inches(5), Inches(0.4), "Propiedades Fisicoquímicas", font_size=18, color=GREEN_DARK, bold=True, font_name='Georgia')
add_table(slide, Inches(1.0), Inches(2.2), Inches(5), fq_data)
org_data = [["Parámetro","Valor"],["Estado","Líquido"],["Color","Verde"],["Olor","Eucalipto"],["Turbidez","Sin turbidez"]]
add_text_box(slide, Inches(7.0), Inches(1.7), Inches(5), Inches(0.4), "Propiedades Organolépticas", font_size=18, color=GREEN_DARK, bold=True, font_name='Georgia')
add_table(slide, Inches(7.0), Inches(2.2), Inches(5), org_data)

# ════════════════════════════════════════════════════════════════
# SLIDE 19: ENVASE
# ════════════════════════════════════════════════════════════════
slide = dark_content_slide("Envase y Rotulación", "Producto Elaborado")
env_data = [["Propiedad","Detalle"],["Material","PET virgen"],["Capacidad","100 mL"],["Tipo","Atomizador spray"],["Reciclable","Sí — SPI ♻"],["Rotulación","Según ANMAT"]]
add_table(slide, Inches(1.0), Inches(2.0), Inches(5), env_data, header_color=RGBColor(0x22,0x44,0x33))

add_text_box(slide, Inches(7.0), Inches(1.7), Inches(5), Inches(0.4), "Embalaje y Distribución", font_size=18, color=GREEN_PALE, bold=True, font_name='Georgia')
emb_items = [
    ("🧴", "Botella 100 mL", "Envase primario con atomizador spray"),
    ("📦", "Caja Cartón Corrugado", "Embalaje secundario para agrupación"),
    ("🏗️", "Pallet Retornable", "Plástico reutilizable + film stretch"),
]
for i, (icon, title, desc) in enumerate(emb_items):
    y = Inches(2.2) + i * Inches(1.5)
    add_card(slide, Inches(7.0), y, Inches(5.2), Inches(1.2), title, desc,
             bg_color=RGBColor(0x22,0x44,0x33), title_color=GREEN_PALE,
             text_color=RGBColor(0xAA,0xBB,0xAA), icon=icon)

# ════════════════════════════════════════════════════════════════
# SLIDE 20: SEP LOCALIZACION
# ════════════════════════════════════════════════════════════════
separator_slide("Localización", "Selección macro y micro de la planta", "📍")

# ════════════════════════════════════════════════════════════════
# SLIDE 21: MACROLOCALIZACION
# ════════════════════════════════════════════════════════════════
slide = cream_content_slide("Macrolocalización: Santa Fe", "Localización")
macro_data = [
    ["Factor", "Santa Fe", "Buenos Aires"],
    ["Materia prima", "Alta", "Media"],
    ["Mano de obra", "Disponible", "Disponible"],
    ["Transporte", "RN 11, RN 19", "Saturado"],
    ["Parques industriales", "Disponibles", "Saturados"],
    ["Energía", "Disponible", "Disponible"],
]
add_table(slide, Inches(1.0), Inches(1.9), Inches(7), macro_data)

shape21 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(5.0), Inches(7), Inches(1.0))
shape21.fill.solid(); shape21.fill.fore_color.rgb = GREEN_DARK; shape21.line.fill.background()
tf21 = shape21.text_frame; tf21.margin_left = Inches(0.2); tf21.word_wrap = True
p21 = tf21.paragraphs[0]
r21 = p21.add_run(); r21.text = "✓ Santa Fe seleccionada\n"; r21.font.size = Pt(16); r21.font.color.rgb = WHITE; r21.font.bold = True
r21b = p21.add_run(); r21b.text = "Producción local de eucalipto, infraestructura vial (RN 11, RN 19), 2 aeropuertos, parques industriales activos."
r21b.font.size = Pt(12); r21b.font.color.rgb = RGBColor(0xCC,0xDD,0xCC)

# ════════════════════════════════════════════════════════════════
# SLIDE 22: MICROLOCALIZACION
# ════════════════════════════════════════════════════════════════
slide = dark_content_slide("Microlocalización: Sauce Viejo", "Localización")
micro_items = [
    "Proveedor: Forestal Sauce Viejo S.A. (hojas como residuo maderero)",
    "22 km de Santa Fe capital",
    "Parque Industrial con accesos pavimentados",
    "Central eléctrica 280 MW",
    "Servicios de Aguas Santafesinas",
    "Población: ~17.000 hab. — RN 11",
]
add_bullet_list(slide, Inches(1.0), Inches(1.9), Inches(5.5), micro_items, color=RGBColor(0xCC,0xDD,0xCC), font_size=15)

stats22 = [("22 km", "de Santa Fe capital"), ("280 MW", "Central eléctrica"), ("17.000", "Habitantes"), ("RN 11", "Ruta Nacional")]
for i, (num, label) in enumerate(stats22):
    col = i % 2; row = i // 2
    x = Inches(7.5) + col * Inches(2.8)
    y = Inches(1.9) + row * Inches(2.0)
    add_stat_card(slide, x, y, Inches(2.5), Inches(1.6), num, label,
                  bg_color=RGBColor(0x22,0x44,0x33), num_color=GREEN_PALE, label_color=RGBColor(0x88,0xAA,0x88))

# ════════════════════════════════════════════════════════════════
# SLIDE 23: SEP PROCESO
# ════════════════════════════════════════════════════════════════
separator_slide("Proceso Productivo", "9 etapas de producción", "⚙️")

# ════════════════════════════════════════════════════════════════
# SLIDE 24: METODOS
# ════════════════════════════════════════════════════════════════
slide = cream_content_slide("Métodos de Obtención del Aceite Esencial", "Proceso Productivo")
methods = [
    ("♨️", "A. Destilación por arrastre con vapor", "El vapor arrastra aceite volátil. Método tradicional, rendimiento moderado."),
    ("🔧", "B. Extracción por prensado", "Presión mecánica sobre el material vegetal. Bajo rendimiento."),
    ("✅", "C. Extracción con solventes", "Etanol disuelve el aceite. Mayor rendimiento (3%). Método elegido: lixiviación + arrastre con vapor."),
]
for i, (icon, title, text) in enumerate(methods):
    x = Inches(1.0) + i * Inches(3.9)
    border_color = GREEN_DARK if i == 2 else None
    add_card(slide, x, Inches(1.9), Inches(3.5), Inches(2.5), title, text, icon=icon)

# ════════════════════════════════════════════════════════════════
# SLIDE 25: FLUJO
# ════════════════════════════════════════════════════════════════
slide = dark_content_slide("Etapas del Proceso — Flujo Completo", "Proceso Productivo")
etapas = [
    "1. Recepción y control", "2. Separación manual", "3. Lixiviación",
    "4. Arrastre con vapor", "5. Condensación", "6. Decantación",
    "7. Mezclado", "8. Envasado", "9. Almacenamiento"
]
for i, etapa in enumerate(etapas):
    col = i % 5; row = i // 5
    x = Inches(0.5) + col * Inches(2.5)
    y = Inches(2.0) + row * Inches(2.2)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.2), Inches(1.4))
    shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(0x22,0x44,0x33)
    shape.line.color.rgb = RGBColor(0x44,0x66,0x55); shape.line.width = Pt(1)
    tf = shape.text_frame; tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = tf.paragraphs[0].add_run(); r.text = etapa
    r.font.size = Pt(14); r.font.color.rgb = WHITE; r.font.name = 'Calibri'; r.font.bold = True
    # arrows
    if col < 4 and i < 8:
        add_text_box(slide, x + Inches(2.2), y + Inches(0.4), Inches(0.3), Inches(0.5), "→",
                     font_size=18, color=GREEN_PALE, alignment=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDES 26-31: PROCESS DETAILS (condensed)
# ════════════════════════════════════════════════════════════════

# SLIDE 26: RECEPCION
slide = cream_content_slide("Recepción y Separación Manual", "Proceso Productivo")
rec_items = ["Eucalipto en fardos de 25 kg (pallets)","Etanol en bidones de 20 L","Isopropanol en bidones de 20 L","Glicerina en bidones de 5 L","Sistema FIFO"]
add_text_box(slide, Inches(1.0), Inches(1.7), Inches(5), Inches(0.4), "Recepción de Materias Primas", font_size=18, color=GREEN_DARK, bold=True, font_name='Georgia')
add_bullet_list(slide, Inches(1.0), Inches(2.2), Inches(5), rec_items)
add_text_box(slide, Inches(7.0), Inches(1.7), Inches(5), Inches(0.4), "Balance de Separación", font_size=18, color=GREEN_DARK, bold=True, font_name='Georgia')
add_text_box(slide, Inches(7.0), Inches(2.4), Inches(5), Inches(2.5),
             "F₁ = 100 kg/h (entrada)\n\n→ F₂ = 79 kg/h hojas\n→ R₁ = 20 kg/h ramas\n→ Pérdidas = 1%",
             font_size=16, color=TEXT_DARK)

# SLIDE 27: LIXIVIACION
slide = dark_content_slide("Lixiviación — Parámetros", "Proceso Productivo")
lix_data = [["Parámetro","Valor"],["Solvente","Etanol 96%"],["Relación S/L","1:5"],["Tiempo","4 horas"],["T° medio líquido","80 °C"],["Vapor calentamiento","148 °C"],["ΔT","68 °C"],["T° condensación","85 °C"]]
add_table(slide, Inches(1.0), Inches(2.0), Inches(5.5), lix_data, header_color=RGBColor(0x22,0x44,0x33))
add_text_box(slide, Inches(7.5), Inches(1.7), Inches(4), Inches(0.4), "Sistema Integrado", font_size=18, color=GREEN_PALE, bold=True, font_name='Georgia')
equip = [("Percolador", "V = 0,62 m³ · D = 0,92 m"), ("Evaporador", "A = 0,31 m²"), ("Condensador", "A = 0,49 m²")]
for i, (name, spec) in enumerate(equip):
    add_card(slide, Inches(7.5), Inches(2.2) + i * Inches(1.3), Inches(4.5), Inches(1.0), name, spec,
             bg_color=RGBColor(0x22,0x44,0x33), title_color=GREEN_PALE, text_color=RGBColor(0xAA,0xBB,0xAA))

# SLIDE 28: ARRASTRE
slide = cream_content_slide("Arrastre con Vapor y Condensación", "Proceso Productivo")
arr_items = ["Concentra aceite esencial","Vapor inyectado: 3,50 kg/h","78% del vapor concentra el aceite","90% del aceite se recupera","F₆ = 4,26 kg/h"]
add_text_box(slide, Inches(1.0), Inches(1.7), Inches(5), Inches(0.4), "Arrastre con Vapor", font_size=18, color=GREEN_DARK, bold=True, font_name='Georgia')
add_bullet_list(slide, Inches(1.0), Inches(2.2), Inches(5), arr_items)
cond_items = ["Cambio de estado sin pérdida: F₆ = F₇ = 4,26 kg/h","Destilador: D = 0,14 m, H = 0,13 m","T. calentamiento: 10,50 min","T. extracción: 30,43 min","Potencia: 3,21 kW"]
add_text_box(slide, Inches(7.0), Inches(1.7), Inches(5), Inches(0.4), "Condensación y Equipo", font_size=18, color=GREEN_DARK, bold=True, font_name='Georgia')
add_bullet_list(slide, Inches(7.0), Inches(2.2), Inches(5), cond_items)

# SLIDE 29: DECANTACION
slide = dark_content_slide("Decantación y Mezclado", "Proceso Productivo")
dec_items = ["Separación por gravedad","Aceite esencial: 1,38 kg/h","Pérdidas: 10%","Agua residual: 2,88 kg/h"]
add_text_box(slide, Inches(1.0), Inches(1.7), Inches(5), Inches(0.4), "Decantación", font_size=18, color=GREEN_PALE, bold=True, font_name='Georgia')
add_bullet_list(slide, Inches(1.0), Inches(2.2), Inches(5), dec_items, color=RGBColor(0xCC,0xDD,0xCC))
add_text_box(slide, Inches(7.0), Inches(1.7), Inches(5), Inches(0.4), "Mezclado", font_size=18, color=GREEN_PALE, bold=True, font_name='Georgia')
add_text_box(slide, Inches(7.0), Inches(2.3), Inches(5), Inches(1.5),
             "Aceite 1,38 kg/h + Isopropanol 1,44 kg/h + Glicerina 0,02 kg/h\n\n= 2,84 kg/h repelente",
             font_size=16, color=RGBColor(0xCC,0xDD,0xCC))

# SLIDE 30: ENVASADO
slide = cream_content_slide("Envasado y Almacenamiento", "Proceso Productivo")
add_card(slide, Inches(1.0), Inches(1.9), Inches(5.4), Inches(2.0),
         "Máquina ZT-Pack CBZ-4", "495 botellas/día (247/turno)\nPotencia: 3 kW\nBotella PET 100 mL con atomizador spray")
add_text_box(slide, Inches(7.0), Inches(1.7), Inches(5), Inches(0.4), "Almacenamiento", font_size=18, color=GREEN_DARK, bold=True, font_name='Georgia')
alm_items = ["Cajas de cartón corrugado","Pallets retornables de plástico","Film stretch","Sistema FIFO"]
add_bullet_list(slide, Inches(7.0), Inches(2.2), Inches(5), alm_items)

# SLIDE 31: BALANCE
slide = dark_content_slide("Balance de Masa — Resumen", "Proceso Productivo")
bal_data = [
    ["Etapa","Entrada (kg/h)","Salida (kg/h)","Residuo"],
    ["Recepción","100","100","—"],
    ["Separación","100","79","21 (ramas)"],
    ["Lixiviación","471,89","1,70 (AE) + 115,95","354,24 (S)"],
    ["Arrastre","1,70 + 3,50","4,26","0,94"],
    ["Condensación","4,26","4,26","—"],
    ["Decantación","4,26","1,38","2,88 (agua)"],
    ["Mezclado","1,38 + 1,44 + 0,02","2,84","—"],
    ["Envasado","2,84","49,50 L/día = 495 bot","—"],
]
add_table(slide, Inches(0.8), Inches(1.8), Inches(11.5), bal_data, header_color=RGBColor(0x22,0x44,0x33))

# ════════════════════════════════════════════════════════════════
# SLIDE 32: SEP CALIDAD
# ════════════════════════════════════════════════════════════════
separator_slide("Control de Calidad", "Aseguramiento en cada etapa", "🛡️")

# SLIDE 33: QC MP
slide = cream_content_slide("Control de Calidad — Materias Primas", "Control de Calidad")
qc_cards = [
    ("🌿", "Eucalipto", "Organolépticos, morfológicos, fisicoquímicos\n(aceites ≥ 20 mL/g, secado ≤ 10%, cenizas ≤ 6%),\nmicrobiológicos."),
    ("🧪", "Etanol", "Color, olor, densidad, acidez.\nVerificación de concentración 96%."),
    ("⚗️", "Isopropanol", "Color, olor, densidad, acidez.\nControl de pureza y certificación."),
    ("🫧", "Glicerina", "Color, densidad ≥ 1,25, agua ≤ 5%,\nresiduo de ignición ≤ 0,01%."),
]
for i, (icon, title, text) in enumerate(qc_cards):
    col = i % 2; row = i // 2
    x = Inches(1.0) + col * Inches(5.8)
    y = Inches(1.9) + row * Inches(2.4)
    add_card(slide, x, y, Inches(5.4), Inches(2.0), title, text, icon=icon)

# SLIDE 34: QC PROCESO
slide = dark_content_slide("Control en Proceso y Producto Final", "Control de Calidad")
add_text_box(slide, Inches(1.0), Inches(1.7), Inches(5), Inches(0.4), "Aceite Esencial (en proceso)", font_size=18, color=GREEN_PALE, bold=True, font_name='Georgia')
qcp_items = ["Control organoléptico continuo","Cromatografía de capa fina (TLC)","Determinación de componentes por Rf"]
add_bullet_list(slide, Inches(1.0), Inches(2.2), Inches(5), qcp_items, color=RGBColor(0xCC,0xDD,0xCC))

pf_data = [["Parámetro","Especificación"],["Aspecto","Líquido verde"],["Olor","Eucalipto"],["pH","4,50"],["Aceite esencial","Verificado por TLC"]]
add_text_box(slide, Inches(7.0), Inches(1.7), Inches(5), Inches(0.4), "Producto Final", font_size=18, color=GREEN_PALE, bold=True, font_name='Georgia')
add_table(slide, Inches(7.0), Inches(2.2), Inches(5), pf_data, header_color=RGBColor(0x22,0x44,0x33))

# SLIDE 35: HYS
slide = cream_content_slide("Higiene y Seguridad — EPP", "Higiene y Seguridad")
hys_cards = [
    ("🦺", "EPP Obligatorio", "Casco ABS blanco, anteojos 3M, protector auditivo,\nzapatos funcionales, guantes de cuero,\nmascarilla con válvula."),
    ("🏭", "Condiciones Ambientales", "Ventilación controlada, iluminación según norma,\ncontrol de ruidos y vibraciones."),
    ("🧯", "Protección Contra Incendio", "Clasificación de riesgo, extintores adecuados,\nplan de evacuación, señalización reglamentaria."),
]
for i, (icon, title, text) in enumerate(hys_cards):
    x = Inches(1.0) + i * Inches(3.9)
    add_card(slide, x, Inches(1.9), Inches(3.5), Inches(2.5), title, text, icon=icon)

# SLIDE 36: CONDICIONES
slide = dark_content_slide("Condiciones Ambientales e Incendio", "Higiene y Seguridad")
cond_items = ["Ventilación controlada","Iluminación por método de cavidades zonales","Ruidos y vibraciones controlados","Instalaciones eléctricas seguras","Colores en cañerías según DIN 2.403"]
add_text_box(slide, Inches(1.0), Inches(1.7), Inches(5), Inches(0.4), "Condiciones de Trabajo", font_size=18, color=GREEN_PALE, bold=True, font_name='Georgia')
add_bullet_list(slide, Inches(1.0), Inches(2.2), Inches(5), cond_items, color=RGBColor(0xCC,0xDD,0xCC))
inc_items = ["Clasificación de riesgo del sector","Extintores adecuados según tipo de fuego","Plan de evacuación documentado","Señalización de obligación, advertencia e informativa","Rombo NFPA para productos almacenados"]
add_text_box(slide, Inches(7.0), Inches(1.7), Inches(5), Inches(0.4), "Protección Contra Incendio", font_size=18, color=GREEN_PALE, bold=True, font_name='Georgia')
add_bullet_list(slide, Inches(7.0), Inches(2.2), Inches(5), inc_items, color=RGBColor(0xCC,0xDD,0xCC))

# SLIDE 37: IMPACTO AMBIENTAL
slide = cream_content_slide("Impacto Ambiental", "Impacto y Marco Jurídico")
add_text_box(slide, Inches(1.0), Inches(1.7), Inches(5), Inches(0.4), "Metodología: Matriz de Conesa", font_size=18, color=GREEN_DARK, bold=True, font_name='Georgia')
factores = "Signo · Intensidad · Extensión · Momento · Persistencia · Reversibilidad · Recuperabilidad · Sinergia · Acumulación · Efecto · Periodicidad"
add_text_box(slide, Inches(1.0), Inches(2.3), Inches(5.5), Inches(1.5), factores, font_size=13, color=TEXT_MID)
add_card(slide, Inches(7.0), Inches(1.9), Inches(5), Inches(2.0), "Resultados",
         "Impactos compatibles y moderados en etapas de construcción y operación.\n\nPlan de mitigación para ambas etapas del proyecto.")

# SLIDE 38: IMPACTO SOCIAL
slide = dark_content_slide("Impacto Social y Económico", "Impacto y Marco Jurídico")
imp_cards = [
    ("👥", "Social", "20 empleados directos + 7 tercerizados.\nPromoción de alternativas saludables.\nFortalecimiento industrial local."),
    ("💰", "Económico", "Materia prima de bajo costo.\nAlto valor agregado.\nValorización de recursos forestales."),
    ("🌍", "Ambiental", "Producto biodegradable.\nSin químicos sintéticos.\nBajo impacto ecológico."),
]
for i, (icon, title, text) in enumerate(imp_cards):
    x = Inches(1.0) + i * Inches(3.9)
    add_card(slide, x, Inches(1.9), Inches(3.5), Inches(2.5), title, text,
             bg_color=RGBColor(0x22,0x44,0x33), title_color=GREEN_PALE, text_color=RGBColor(0xBB,0xCC,0xBB), icon=icon)

# SLIDE 39: MARCO JURIDICO
slide = cream_content_slide("Marco Jurídico", "Impacto y Marco Jurídico")
mj_cards = [
    ("Registro", "ANMAT — producto y establecimiento"),
    ("Localización", "Ley 8.478 (promoción industrial)\nLey 11.525 (parques industriales)"),
    ("Laboral", "LCT · Convenio colectivo\nEscala salarial CCT 42/89"),
    ("Calidad", "Farmacopea Argentina"),
    ("Fiscal", "Monotributo / Responsable inscripto"),
    ("Seguridad", "Ley 19.587 — Higiene y Seguridad"),
]
for i, (title, text) in enumerate(mj_cards):
    col = i % 3; row = i // 3
    x = Inches(1.0) + col * Inches(3.9)
    y = Inches(1.9) + row * Inches(2.2)
    add_card(slide, x, y, Inches(3.5), Inches(1.8), title, text)

# SLIDE 40: ORGANIZACION
slide = dark_content_slide("Organización Industrial", "Organización")
org_data = [
    ["Sector","Cantidad"],
    ["Gerente General","1"],["Jefe Producción y Calidad","1"],["Jefe Admin. y Comercialización","1"],
    ["Operarios producción","5"],["Mantenimiento","2"],["Analistas calidad","2"],
    ["Administrativos","1"],["Recepcionista","1"],["Total directos","20"],
]
add_text_box(slide, Inches(1.0), Inches(1.7), Inches(5), Inches(0.4), "Estructura — SRL", font_size=18, color=GREEN_PALE, bold=True, font_name='Georgia')
add_table(slide, Inches(1.0), Inches(2.2), Inches(5.5), org_data, header_color=RGBColor(0x22,0x44,0x33))
add_text_box(slide, Inches(7.5), Inches(1.7), Inches(4.5), Inches(0.4), "Servicios Tercerizados (7)", font_size=18, color=GREEN_PALE, bold=True, font_name='Georgia')
terc = ["Limpieza","Higiene y seguridad","Seguridad privada","Legal","Transporte","Laboratorio externo","Informática"]
add_bullet_list(slide, Inches(7.5), Inches(2.2), Inches(4.5), terc, color=RGBColor(0xCC,0xDD,0xCC))

# SLIDE 41: EDIFICACION
slide = cream_content_slide("Planificación y Edificación", "Organización")
edif_data = [["Sector","m²"],["Producción","539,00"],["Almacén MP","91,30"],["Oficinas","119,31"],["Laboratorio","36,07"],["Mantenimiento","52,49"],["Expedición","94,62"],["Baños","99,20"]]
add_text_box(slide, Inches(1.0), Inches(1.7), Inches(5), Inches(0.4), "Zona Cubierta — 1.127,27 m²", font_size=18, color=GREEN_DARK, bold=True, font_name='Georgia')
add_table(slide, Inches(1.0), Inches(2.2), Inches(5), edif_data)
add_stat_card(slide, Inches(7.5), Inches(1.9), Inches(4.5), Inches(1.5), "4.895 m²", "Terreno total",
              bg_color=GREEN_DARK, num_color=WHITE, label_color=RGBColor(0xCC,0xDD,0xCC))
add_text_box(slide, Inches(7.5), Inches(3.8), Inches(4.5), Inches(0.4), "Zona Descubierta — 3.767,73 m²", font_size=18, color=GREEN_DARK, bold=True, font_name='Georgia')
desc_items = ["Calles internas","Estacionamiento","Espacios verdes","Senderos peatonales"]
add_bullet_list(slide, Inches(7.5), Inches(4.3), Inches(4.5), desc_items)

# SLIDE 42: SERVICIOS AGUA Y GAS
slide = dark_content_slide("Servicios Auxiliares — Agua y Gas", "Servicios Auxiliares")
agua_data = [["Uso","m³/sem"],["Vapor","20,59"],["Refrigeración","12,25"],["Limpieza","7,22"],["Uso humano","2,70"]]
add_text_box(slide, Inches(1.0), Inches(1.7), Inches(5), Inches(0.4), "💧 Agua Total: 44,80 m³/sem", font_size=18, color=GREEN_PALE, bold=True, font_name='Georgia')
add_table(slide, Inches(1.0), Inches(2.2), Inches(5), agua_data, header_color=RGBColor(0x22,0x44,0x33))
gas_data = [["Uso","m³/día"],["Lixiviación","6,45"],["Arrastre","3,82"],["Calefacción","28,38"]]
add_text_box(slide, Inches(7.0), Inches(1.7), Inches(5), Inches(0.4), "🔥 Gas Natural: 161 m³/sem", font_size=18, color=GREEN_PALE, bold=True, font_name='Georgia')
add_table(slide, Inches(7.0), Inches(2.2), Inches(5), gas_data, header_color=RGBColor(0x22,0x44,0x33))

# SLIDE 43: ELECTRICIDAD
slide = cream_content_slide("Electricidad, Aire y Caldera", "Servicios Auxiliares")
elec_data = [["Uso","kWh/sem"],["Motores","953,80"],["Iluminación","1.094,05"]]
add_text_box(slide, Inches(1.0), Inches(1.7), Inches(5), Inches(0.4), "⚡ Energía: 2.047,85 kWh/sem", font_size=18, color=GREEN_DARK, bold=True, font_name='Georgia')
add_table(slide, Inches(1.0), Inches(2.2), Inches(5), elec_data)
add_card(slide, Inches(1.0), Inches(4.0), Inches(5), Inches(1.0), "💨 Aire Comprimido", "200 L/min = 960.000 L/sem")
cald_data = [["Parámetro","Valor"],["Vapor","218,73 kg/h"],["Agua","257,4 kg/h"],["Eficiencia","85%"]]
add_text_box(slide, Inches(7.0), Inches(1.7), Inches(5), Inches(0.4), "Caldera LNR 15", font_size=18, color=GREEN_DARK, bold=True, font_name='Georgia')
add_table(slide, Inches(7.0), Inches(2.2), Inches(5), cald_data)

# ════════════════════════════════════════════════════════════════
# SLIDE 44: SEP INVERSIONES
# ════════════════════════════════════════════════════════════════
separator_slide("Inversiones y Costos", "Análisis económico del proyecto", "💵")

# SLIDE 45: INVERSION
slide = dark_content_slide("Inversión Total", "Inversiones y Costos")
inv_data = [
    ["Rubro","Detalle","Costo (miles $)"],
    ["I","Terreno y edificios","289.996,60"],
    ["II","Equipos y accesorios","608.108,23"],
    ["III","Instalaciones eléctricas","3.599,60"],
    ["IV","Equipamiento oficinas","14.745,00"],
    ["V","Rodados","67.720,00"],
    ["VI","Gastos de proyecto","54,10"],
    ["VII","Montaje","35,92"],
    ["TOTAL","Con 5% protección","$1.033.472,43k"],
]
add_table(slide, Inches(0.8), Inches(1.8), Inches(8), inv_data, header_color=RGBColor(0x22,0x44,0x33))

shape45 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.5), Inches(1.9), Inches(3.3), Inches(1.5))
shape45.fill.solid(); shape45.fill.fore_color.rgb = RGBColor(0x1F,0x35,0x28); shape45.line.fill.background()
tf45 = shape45.text_frame; tf45.margin_left = Inches(0.2); tf45.word_wrap = True
p45a = tf45.paragraphs[0]; r45a = p45a.add_run(); r45a.text = "FINANCIAMIENTO"; r45a.font.size = Pt(9); r45a.font.color.rgb = RGBColor(0x88,0xAA,0x88); r45a.font.bold = True
p45b = tf45.add_paragraph(); r45b = p45b.add_run(); r45b.text = "Banco Nación Argentina\n25% del total · Tasa 40%"
r45b.font.size = Pt(13); r45b.font.color.rgb = WHITE; r45b.font.bold = True

# SLIDE 46: COSTOS
slide = cream_content_slide("Costos Operativos", "Inversiones y Costos")
cost_items = ["Costo materias primas (eucalipto + solventes + glicerina)","Envase y embalaje","Mano de obra directa e indirecta","Energía eléctrica y combustible (gas)","Agua, EPP, mantenimiento","Seguros y servicios tercerizados","Depreciaciones y amortizaciones"]
add_text_box(slide, Inches(1.0), Inches(1.7), Inches(5), Inches(0.4), "Componentes del COT", font_size=18, color=GREEN_DARK, bold=True, font_name='Georgia')
add_bullet_list(slide, Inches(1.0), Inches(2.2), Inches(5.5), cost_items)
add_stat_card(slide, Inches(7.5), Inches(2.0), Inches(4.5), Inches(1.8), "$33.421", "Costo unitario / botella",
              bg_color=GREEN_DARK, num_color=WHITE, label_color=RGBColor(0xCC,0xDD,0xCC))
add_stat_card(slide, Inches(7.5), Inches(4.2), Inches(4.5), Inches(1.8), "$41.776", "Precio de venta (25% ganancia)",
              bg_color=GREEN_DARK, num_color=GOLD_LIGHT, label_color=RGBColor(0xCC,0xDD,0xCC))

# SLIDE 47: PUNTO EQUILIBRIO
slide = dark_content_slide("Punto de Equilibrio", "Evaluación")
stats47 = [("73.677", "Botellas — Punto de equilibrio"), ("49,6%", "de la capacidad instalada"), ("43,1%", "Margen de seguridad")]
for i, (num, label) in enumerate(stats47):
    y = Inches(2.0) + i * Inches(1.7)
    add_stat_card(slide, Inches(1.0), y, Inches(5), Inches(1.4), num, label,
                  bg_color=RGBColor(0x22,0x44,0x33), num_color=GREEN_PALE, label_color=RGBColor(0x88,0xAA,0x88))
add_text_box(slide, Inches(7.0), Inches(1.7), Inches(5), Inches(0.4), "Análisis", font_size=18, color=GREEN_PALE, bold=True, font_name='Georgia')
eq_items = ["Ingreso = PV × Q","Costo = CF + CV × Q","Todo volumen superior al punto de equilibrio genera ganancias","Capacidad real (148.500 bot/año) supera ampliamente el equilibrio"]
add_bullet_list(slide, Inches(7.0), Inches(2.2), Inches(5), eq_items, color=RGBColor(0xCC,0xDD,0xCC))

# SLIDE 48: VAN TIR
slide = cream_content_slide("Evaluación del Proyecto — VAN y TIR", "Evaluación")
van_data = [
    ["Escenario","TMAR","VAN (miles $)","TIR"],
    ["Sin inflación, sin financiamiento","15%","5.571.871","190,28%"],
    ["Con inflación, sin financiamiento","55%","5.571.871","291,88%"],
    ["Con inflación, con financiamiento","52%","6.128.700","381,98%"],
]
add_table(slide, Inches(0.8), Inches(1.8), Inches(11.5), van_data)
add_stat_card(slide, Inches(1.0), Inches(4.5), Inches(5.3), Inches(1.8), "$6.128 M", "VAN — Valor Actual Neto",
              bg_color=GREEN_DARK, num_color=WHITE, label_color=RGBColor(0xCC,0xDD,0xCC))
add_stat_card(slide, Inches(6.8), Inches(4.5), Inches(5.3), Inches(1.8), "381,98%", "TIR — Tasa Interna de Retorno",
              bg_color=GREEN_DARK, num_color=GOLD_LIGHT, label_color=RGBColor(0xCC,0xDD,0xCC))

# SLIDE 49: FACTIBILIDAD
slide = dark_content_slide("Factibilidad del Proyecto", "Evaluación")
fact_items = ["VAN positivo en los 3 escenarios → proyecto viable","TIR muy superior a TMAR → inversión favorable","FNE crecientes en los 5 años de evaluación","Valor de salvamento: $2.379 millones","Escenario 3 elegido (mayor VAN y TIR)","Un turno adicional incrementa rentabilidad"]
add_bullet_list(slide, Inches(1.0), Inches(1.9), Inches(5.5), fact_items, color=RGBColor(0xCC,0xDD,0xCC), font_size=16)

shape49 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.5), Inches(2.5), Inches(4.5), Inches(1.5))
shape49.fill.solid(); shape49.fill.fore_color.rgb = GREEN_DARK; shape49.line.fill.background()
tf49 = shape49.text_frame; tf49.word_wrap = True
p49 = tf49.paragraphs[0]; p49.alignment = PP_ALIGN.CENTER
r49 = p49.add_run(); r49.text = "✓ PROYECTO VIABLE"; r49.font.size = Pt(24); r49.font.color.rgb = WHITE; r49.font.bold = True

add_text_box(slide, Inches(7.5), Inches(4.5), Inches(4.5), Inches(0.8),
             "Precio mínimo rentable:\n$29.190,95/botella",
             font_size=14, color=RGBColor(0x99,0xBB,0x99), alignment=PP_ALIGN.CENTER)

# SLIDE 50: CONCLUSIONES
slide = dark_content_slide("Conclusiones Generales", "Conclusiones")
conclusiones = [
    "Condiciones favorables de implementación técnica y territorial",
    "Proyección de ingresos justifica la inversión (VAN $6.128M, TIR 381,98%)",
    "Aporte al desarrollo local: empleo e industria en Sauce Viejo",
    "Criterios de calidad en cada etapa (Farmacopea Argentina)",
    "Valor agregado a recursos forestales locales",
    "Alternativa sostenible y competitiva en el mercado",
]
add_numbered_list(slide, Inches(1.0), Inches(1.9), Inches(10), conclusiones,
                  color=RGBColor(0xCC,0xDD,0xCC), num_color=GREEN_PALE, font_size=17)

# SLIDE 51: DASHBOARD
slide = cream_content_slide("Indicadores Clave del Proyecto", "Resumen")
dashboard = [
    ("49,50 L/d","Producción diaria"), ("148.500","Botellas / año"),
    ("$1.033 M","Inversión total"), ("$33.421","Costo unitario"),
    ("$41.776","Precio de venta"), ("73.677","Pto. equilibrio"),
    ("$6.128 M","VAN"), ("381,98%","TIR"),
    ("27","Empleados"), ("4.895 m²","Terreno"),
]
for i, (num, label) in enumerate(dashboard):
    col = i % 5; row = i // 5
    x = Inches(0.5) + col * Inches(2.46)
    y = Inches(1.9) + row * Inches(2.3)
    nc = GREEN_DARK if i in [6, 7] else TEXT_DARK
    add_stat_card(slide, x, y, Inches(2.2), Inches(1.9), num, label, num_color=nc)

# ════════════════════════════════════════════════════════════════
# SLIDE 52: GRACIAS
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, GREEN_DEEP)
add_text_box(slide, Inches(0), Inches(1.5), W, Inches(0.6), "🍃",
             font_size=48, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(0), Inches(2.3), W, Inches(1.2), "Gracias",
             font_size=72, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER, font_name='Georgia')
line52 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.8), Inches(3.7), Inches(1.7), Inches(0.02))
line52.fill.solid(); line52.fill.fore_color.rgb = GOLD; line52.line.fill.background()
add_text_box(slide, Inches(0), Inches(4.0), W, Inches(0.4),
             "BRESSO, Romina  ·  MACAGNO, Micaela",
             font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(0), Inches(4.6), W, Inches(0.3),
             "Proyecto Final · Ingeniería Química",
             font_size=13, color=RGBColor(0x99,0x99,0x99), italic=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(0), Inches(5.0), W, Inches(0.3),
             "Universidad Tecnológica Nacional — Facultad Regional San Francisco",
             font_size=12, color=RGBColor(0x99,0x99,0x99), alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(0), Inches(5.5), W, Inches(0.3), "2025",
             font_size=12, color=RGBColor(0x77,0x77,0x77), alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# SAVE
# ════════════════════════════════════════════════════════════════
output_path = "/Users/mauriciobuffa/Desktop/Romi/presentation-claude/Presentacion_Repelente_Natural.pptx"
prs.save(output_path)
print(f"✅ PPTX saved to: {output_path}")
print(f"   Total slides: {len(prs.slides)}")
